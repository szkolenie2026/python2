from pptx import Presentation

# Tworzymy nową prezentację
prs = Presentation()

# Dodajemy slajd tytułowy (layout nr 0)
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Programowanie w Python"
subtitle.text = "Rafał Proskura"

# Dodajemy slajd z tekstem (layout nr 1)
bullet_slide_layout = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(bullet_slide_layout)
shapes = slide2.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Dlaczego Python?"
tf = body_shape.text_frame
tf.text = "Wybrały go giganty: NASA, Google, Youtube"
p = tf.add_paragraph()
p.text = "Najpopularniejszy język na świecie"

# Zapisujemy gotowy plik!
prs.save('Moja_Gotowa_Prezentacja.pptx')
print("Plik .pptx został wygenerowany!")